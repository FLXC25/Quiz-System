import io
import os
import random
import re
from collections import Counter
from typing import Iterable, List, Optional

from flask import (
    Flask,
    flash,
    redirect,
    render_template,
    request,
    session,
    url_for,
)
from PyPDF2 import PdfReader
from pptx import Presentation
from werkzeug.utils import secure_filename

app = Flask(__name__)
app.secret_key = os.environ.get("FLASK_SECRET_KEY", "dev-secret")

ALLOWED_EXTENSIONS = {".pdf", ".ppt", ".pptx"}
DEFAULT_CHOICES = [
    "Definitely true",
    "Likely true",
    "Unlikely",
    "Definitely false",
]


def clean_text(text: str) -> str:
    """Normalize whitespace and drop control characters."""
    return re.sub(r"\s+", " ", text or "").strip()


def extract_text_from_pdf(stream: io.BytesIO) -> str:
    reader = PdfReader(stream)
    text_parts: List[str] = []
    for page in reader.pages:
        extracted = page.extract_text() or ""
        text_parts.append(extracted)
    return "\n".join(text_parts)


def extract_text_from_ppt(stream: io.BytesIO) -> str:
    presentation = Presentation(stream)
    text_parts: List[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_parts.append(shape.text)
    return "\n".join(text_parts)


def extract_text_from_file(file_storage) -> Optional[str]:
    if not file_storage:
        return None

    filename = secure_filename(file_storage.filename or "")
    if not filename:
        return None

    _, ext = os.path.splitext(filename)
    ext = ext.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError("Unsupported file type. Please upload a PDF or PPT/PPTX file.")

    data = file_storage.read()
    stream = io.BytesIO(data)

    if ext == ".pdf":
        return extract_text_from_pdf(stream)
    return extract_text_from_ppt(stream)


def tokenize(text: str) -> List[str]:
    return re.findall(r"[A-Za-z][A-Za-z']+", text)


def build_choices(correct: str, vocabulary: Iterable[str]) -> List[str]:
    distractors = [word for word in vocabulary if word.lower() != correct.lower()]
    random.shuffle(distractors)
    choices = distractors[:3] + [correct]
    while len(choices) < 4:
        choices.append(random.choice(DEFAULT_CHOICES))
    random.shuffle(choices)
    return choices


def generate_question(sentence: str, correct_word: str, vocabulary: Iterable[str]):
    pattern = re.compile(re.escape(correct_word), re.IGNORECASE)

    def replacer(match):
        return "_____"

    question_text = pattern.sub(replacer, sentence, count=1)
    choices = build_choices(correct_word, vocabulary)
    answer_index = choices.index(correct_word)
    return {
        "question": question_text,
        "choices": choices,
        "answer_index": answer_index,
    }


def generate_quiz_from_text(text: str, num_questions: int):
    cleaned = clean_text(text)
    if not cleaned:
        raise ValueError("Please provide study material text or upload a supported file.")

    sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", cleaned) if len(s.split()) >= 4]
    if not sentences:
        sentences = [cleaned]

    vocabulary = tokenize(cleaned)
    if not vocabulary:
        vocabulary = DEFAULT_CHOICES

    common_words = [word for word, _ in Counter(vocabulary).most_common()]
    questions = []

    for sentence in sentences:
        if len(questions) >= num_questions:
            break
        candidates = [word for word in tokenize(sentence) if len(word) > 3]
        used_answers = {q["answer"] for q in questions if "answer" in q}
        candidates = [c for c in candidates if c.lower() not in {u.lower() for u in used_answers}]
        if not candidates:
            continue
        correct_word = random.choice(candidates)
        question = generate_question(sentence, correct_word, common_words)
        question["answer"] = correct_word
        questions.append(question)

    # Fallback to templated questions if we don't have enough natural questions
    while len(questions) < num_questions:
        filler_word = random.choice(common_words)
        question_text = f"Which word appeared in the study material?"
        choices = build_choices(filler_word, common_words)
        answer_index = choices.index(filler_word)
        questions.append(
            {
                "question": question_text,
                "choices": choices,
                "answer_index": answer_index,
                "answer": filler_word,
            }
        )

    for q in questions:
        q.pop("answer", None)

    return {
        "source": cleaned,
        "questions": questions,
    }


@app.route("/")
def index():
    return render_template("page1_input.html")


@app.route("/generate", methods=["POST"])
def generate():
    num_raw = request.form.get("num_questions", "5")
    try:
        num_questions = int(num_raw)
    except ValueError:
        num_questions = 5
    num_questions = max(1, min(num_questions, 10))

    source_text = request.form.get("source_text", "")
    upload = request.files.get("source_file")

    try:
        extracted_text = extract_text_from_file(upload)
    except ValueError as exc:
        flash(str(exc))
        return render_template("page1_input.html")

    material = extracted_text or source_text

    try:
        quiz = generate_quiz_from_text(material, num_questions)
    except ValueError as exc:
        flash(str(exc))
        return render_template("page1_input.html")

    session["quiz"] = quiz
    session.pop("results", None)
    session.modified = True
    return redirect(url_for("show_quiz"))


@app.route("/quiz")
def show_quiz():
    quiz = session.get("quiz")
    if not quiz:
        flash("Generate a quiz first by providing study material.")
        return redirect(url_for("index"))
    return render_template("page2_quiz.html", quiz=quiz)


@app.route("/generate", methods=["GET"])
def redirect_to_index():
    return redirect(url_for("index"))


@app.route("/submit", methods=["POST"])
def submit():
    quiz = session.get("quiz")
    if not quiz:
        flash("Your session expired. Please generate a new quiz.")
        return redirect(url_for("index"))

    user_answers = []
    correct = 0
    for idx, question in enumerate(quiz["questions"]):
        selected = request.form.get(f"q_{idx}")
        try:
            selected_index = int(selected)
        except (TypeError, ValueError):
            selected_index = -1
        user_answers.append(selected_index)
        if selected_index == question["answer_index"]:
            correct += 1

    total = len(quiz["questions"]) or 1
    score = int(round(100 * correct / total))
    session["results"] = {
        "quiz": quiz,
        "user_answers": user_answers,
        "correct": correct,
        "total": total,
        "score": score,
    }
    return redirect(url_for("show_results"))


@app.route("/results")
def show_results():
    results = session.get("results")
    if not results:
        flash("Take the quiz to see your results.")
        return redirect(url_for("index"))
    return render_template(
        "page3_result.html",
        quiz=results["quiz"],
        user_answers=results["user_answers"],
        correct=results["correct"],
        total=results["total"],
        score=results["score"],
    )


if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=int(os.environ.get("PORT", 5000)))
